import os
import re
import tempfile
from datetime import datetime
# 多格式文件解析依赖（保留PyPDF2）
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
# OCR相关依赖（处理图片扫描件）
import pytesseract
from PIL import Image
import fitz  # PyMuPDF，提取PDF中的图片（云端兼容）

# -------------------------- 全局配置（动态获取项目内Tesseract路径，云端友好） --------------------------
# 1. 动态计算项目根目录（无论部署到哪里，自动适配）
BASE_PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# 2. Tesseract OCR路径（指向项目内的tools目录）
TESSERACT_PATH = os.path.join(BASE_PROJECT_DIR, "tools",  "tesseract.exe")
# 3. 语言包路径（指向项目内的tessdata，避免云端找不到语言包）
TESSDATA_DIR = os.path.join(BASE_PROJECT_DIR, "tools", "tesseract-ocr", "tessdata")
# 4. 配置pytesseract（优先用项目内的Tesseract）
pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
os.environ["TESSDATA_PREFIX"] = TESSDATA_DIR  # 指定语言包目录

# 5. 其他全局配置（云端内存保护）
MEDIA_DIR = os.path.join(BASE_PROJECT_DIR, "media")
SUPPORTED_EXTENSIONS = (".txt", ".pdf", ".docx", ".xlsx", ".ppt", ".pptx")
global_knowledge_base = {
    "documents": [],  # 文本块+来源
    "file_list": []  # 已处理文件列表
}
# 大文件适配（云端内存友好）
MIN_TEXT_PER_FILE = 20000
MAX_TEXT_PER_FILE = 100000
MAX_CHUNKS_PER_FILE = 50
FILE_SIZE_THRESHOLD1 = 10 * 1024 * 1024  # 10MB
FILE_SIZE_THRESHOLD2 = 50 * 1024 * 1024  # 50MB


# -------------------------- 工具函数：动态计算提取长度 --------------------------
def get_max_text_length(file_size):
    """根据文件大小动态调整提取文本长度（云端内存保护）"""
    if file_size <= FILE_SIZE_THRESHOLD1:
        return max(MIN_TEXT_PER_FILE, file_size // 2)
    elif file_size <= FILE_SIZE_THRESHOLD2:
        ratio = (file_size - FILE_SIZE_THRESHOLD1) / (FILE_SIZE_THRESHOLD2 - FILE_SIZE_THRESHOLD1)
        return int(MIN_TEXT_PER_FILE + (MAX_TEXT_PER_FILE - MIN_TEXT_PER_FILE) * ratio)
    else:
        return MAX_TEXT_PER_FILE


# -------------------------- 工具函数：文本分块 --------------------------
def split_text_into_chunks(text):
    """文本分块（保留语义，云端内存友好）"""
    chunks = re.split(r'[。！？；\n]{1,2}', text)
    valid_chunks = []
    for chunk in chunks:
        chunk = chunk.strip()
        # 过滤无效文本，保留有效内容
        if chunk and len(chunk) > 15 and len(chunk) < 800 and re.search(r'[一-龥a-zA-Z]', chunk):
            valid_chunks.append(chunk)
    # 限制单文件最大块数，避免云端内存溢出
    return valid_chunks[:MAX_CHUNKS_PER_FILE]


# -------------------------- 核心OCR函数：处理图片扫描件PDF --------------------------
def ocr_pdf_scanned(file_path):
    """
    调用项目内的Tesseract OCR处理图片扫描件PDF
    流程：提取PDF中的图片 → 临时保存 → OCR识别 → 删除临时文件
    """
    text = ""
    try:
        # 1. 用PyMuPDF打开PDF，提取每页图片（云端兼容）
        pdf = fitz.open(file_path)
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            image_list = page.get_images(full=True)
            if not image_list:
                continue

            # 2. 逐个处理图片（避免一次性加载所有图片占内存）
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]

                # 3. 临时保存图片（云端临时目录兼容）
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
                    temp_img.write(image_bytes)
                    temp_img_path = temp_img.name

                # 4. OCR识别（指定简体中文，项目内语言包）
                img = Image.open(temp_img_path)
                # 优化识别：灰度化+二值化（提升扫描件识别率）
                img = img.convert('L')  # 灰度化
                img = img.point(lambda x: 0 if x < 128 else 255, '1')  # 二值化
                ocr_text = pytesseract.image_to_string(img, lang="chi_sim")  # chi_sim=简体中文

                # 5. 拼接文本，控制长度
                text += ocr_text + "\n"
                if len(text) >= MAX_TEXT_PER_FILE:
                    text = text[:MAX_TEXT_PER_FILE]
                    break

                # 6. 立即删除临时文件（云端磁盘友好）
                os.unlink(temp_img_path)

            if len(text) >= MAX_TEXT_PER_FILE:
                break
        pdf.close()
        return text.strip()

    except FileNotFoundError as e:
        print(f"❌ Tesseract OCR未找到：{e}，请检查tools/tesseract-ocr目录是否包含tesseract.exe")
        return ""
    except Exception as e:
        print(f"❌ OCR识别失败：{str(e)}")
        return ""


# -------------------------- 核心函数：解析单个文件 --------------------------
def parse_single_file(file_path):
    """
    解析单个文件：优先PyPDF2提取文本 → 失败则触发OCR（扫描件）
    兼容所有格式+项目内Tesseract+云端部署
    """
    text = ""
    file_ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.basename(file_path)
    file_size = os.path.getsize(file_path)
    max_text_len = get_max_text_length(file_size)

    try:
        print(f"📄 解析文件：{file_name}（大小：{file_size / 1024 / 1024:.1f}MB，计划提取{max_text_len}字符）")

        # TXT文件：直接读取
        if file_ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()[:max_text_len]

        # PDF文件：优先常规提取，失败则OCR
        elif file_ext == ".pdf":
            # 第一步：PyPDF2提取常规文本
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"
                    if len(text) >= max_text_len:
                        text = text[:max_text_len]
                        break

            # 过滤无效文本（仅含空格/符号/数字）
            text_clean = re.sub(r'\s+', ' ', text).strip()
            # 第二步：常规提取失败（<10个有效字符），触发OCR
            if len(text_clean) < 10:
                print(f"⚠️ {file_name} 无文本层，调用项目内Tesseract OCR识别扫描件...")
                ocr_text = ocr_pdf_scanned(file_path)
                if ocr_text:
                    text = ocr_text[:max_text_len]
                else:
                    print(f"❌ {file_name} OCR识别失败：可能是纯图片且无法识别")

        # DOCX文件：读取段落
        elif file_ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    text += para_text + "\n"
                    if len(text) >= max_text_len:
                        text = text[:max_text_len]
                        break

        # Excel文件：读取工作表+行
        elif file_ext == ".xlsx":
            wb = openpyxl.load_workbook(file_path, read_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text += f"=== 工作表：{sheet_name} ===\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None and str(cell).strip()])
                    if row_text:
                        text += row_text + "\n"
                        if len(text) >= max_text_len:
                            text = text[:max_text_len]
                            break
                if len(text) >= max_text_len:
                    break
            wb.close()

        # PPT/PPTX文件：读取幻灯片文本
        elif file_ext in (".ppt", ".pptx"):
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text = shape.text.strip()
                        if slide_text:
                            text += slide_text + "\n"
                            if len(text) >= max_text_len:
                                text = text[:max_text_len]
                                break
                if len(text) >= max_text_len:
                    break

        # 过滤最终无效文本
        text = re.sub(r'\s+', ' ', text).strip()
        if not re.search(r'[一-龥a-zA-Z]', text):
            print(f"⚠️ {file_name} 未提取到有效文本")
            return []

        # 分块并返回
        chunks = split_text_into_chunks(text)
        print(f"✅ 解析完成：{file_name} 提取{len(text)}字符，生成{len(chunks)}个文本块")
        return [{"content": chunk, "source": file_name} for chunk in chunks]

    except Exception as e:
        print(f"❌ 解析文件 {file_name} 失败：{str(e)}")
        return []


# -------------------------- 启动时全量构建知识库 --------------------------
def build_full_kb_on_startup():
    """启动时扫描media（含嵌套文件夹），构建知识库（云端兼容）"""
    global global_knowledge_base
    print(f"🔄 启动时构建知识库：扫描 {MEDIA_DIR}（含嵌套文件夹）...")

    total_files = 0
    success_files = 0
    total_chunks = 0

    # 递归扫描所有嵌套文件夹
    for root, dirs, files in os.walk(MEDIA_DIR):
        for file in files:
            file_path = os.path.join(root, file)
            file_ext = os.path.splitext(file)[1].lower()

            # 仅处理支持的格式+未处理过的文件
            if file_ext in SUPPORTED_EXTENSIONS and file not in global_knowledge_base["file_list"]:
                total_files += 1
                chunks = parse_single_file(file_path)
                if chunks:
                    global_knowledge_base["documents"].extend(chunks)
                    global_knowledge_base["file_list"].append(file)
                    success_files += 1
                    total_chunks += len(chunks)

    print(f"✅ 全量知识库构建完成：扫描{total_files}个文件，成功{success_files}个，生成{total_chunks}个文本块")


# -------------------------- 增量更新知识库 --------------------------
def add_single_file_to_kb(file_path):
    """增量添加单个文件到知识库（管理员上传时触发）"""
    global global_knowledge_base
    file_name = os.path.basename(file_path)

    # 跳过已处理的文件
    if file_name in global_knowledge_base["file_list"]:
        print(f"⚠️ {file_name} 已在知识库中，无需重复添加")
        return True

    # 解析并添加
    chunks = parse_single_file(file_path)
    if not chunks:
        print(f"❌ 增量添加失败：{file_name} 无有效内容或解析失败")
        return False

    global_knowledge_base["documents"].extend(chunks)
    global_knowledge_base["file_list"].append(file_name)
    print(f"✅ 增量添加成功：{file_name}，新增{len(chunks)}个文本块")
    return True


# -------------------------- 关键词检索（适配云端） --------------------------
def retrieve_knowledge(query, top_k=3):  # 默认top_k=3
    global global_knowledge_base
    if not global_knowledge_base["documents"]:
        return []

    stop_words = ["的", "了", "是", "在", "有", "我", "你", "他", "这", "那", "吗", "呢", "就", "都", "也", "还", "和",
                  "及", "与"]
    query_words = [word for word in re.findall(r'\w+', query) if word not in stop_words and len(word) >= 2]
    if not query_words:
        return []

    matched_docs = []
    for doc in global_knowledge_base["documents"]:
        match_count = 0
        word_freq = {}
        # 优化：匹配到1个关键词就计数，无需统计所有（减少循环）
        for word in query_words:
            if word in doc["content"]:
                match_count += 1
                word_freq[word] = doc["content"].count(word)
            # 提前终止：匹配到2个关键词就够了，不用继续
            if match_count >= 2:
                break
        if match_count > 0:
            score = (match_count * 2) + sum(word_freq.values())
            matched_docs.append({
                "content": doc["content"],
                "source": doc["source"],
                "match_score": score
            })
        # 提前终止：收集到top_k*2条就够了，不用遍历所有文档
        if len(matched_docs) >= top_k * 2:
            break

    matched_docs.sort(key=lambda x: x["match_score"], reverse=True)
    return matched_docs[:top_k]